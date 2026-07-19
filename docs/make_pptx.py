from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── 색상 팔레트 ──────────────────────────────────────────────
NAVY   = RGBColor(0x1a, 0x1a, 0x2e)
BLUE   = RGBColor(0x0f, 0x34, 0x60)
ACCENT = RGBColor(0xe9, 0x45, 0x60)
LIGHT  = RGBColor(0xf4, 0xf5, 0xf7)
WHITE  = RGBColor(0xff, 0xff, 0xff)
GRAY   = RGBColor(0x77, 0x77, 0x77)
LGRAY  = RGBColor(0xcc, 0xcc, 0xcc)
GREEN  = RGBColor(0x2b, 0x93, 0x48)
GOLD   = RGBColor(0xf5, 0xa6, 0x23)
PURPLE = RGBColor(0x53, 0x34, 0x83)
TEAL   = RGBColor(0x00, 0x7b, 0x83)

BLANK  = prs.slide_layouts[6]

# ── 헬퍼 ─────────────────────────────────────────────────────
def rect(slide, l, t, w, h, fill, border=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border; s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s

def tx(slide, text, l, t, w, h,
       size=14, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, italic=False, wrap=True):
    b = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = b.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return b

def header(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.33, 1.1, NAVY)
    tx(slide, title, 0.4, 0.13, 12, 0.72, size=26, bold=True, color=WHITE)
    if subtitle:
        tx(slide, subtitle, 0.4, 0.73, 12, 0.34, size=12,
           color=RGBColor(0xaa,0xbb,0xcc), italic=True)
    rect(slide, 0, 1.1, 13.33, 0.04, ACCENT)

def badge(slide, text, l, t, w, h, bg=BLUE, fg=WHITE, size=11):
    rect(slide, l, t, w, h, bg)
    tx(slide, text, l, t+0.04, w, h-0.04, size=size, color=fg,
       align=PP_ALIGN.CENTER, bold=True)

def bullet_lines(slide, items, l, t, w, size=12, color=NAVY, gap=0.42):
    """items: list of (indent, text)"""
    b = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(6))
    tf = b.text_frame; tf.word_wrap = True
    first = True
    for indent, text in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run(); r.text = ("  " * indent) + text
        r.font.size = Pt(size); r.font.color.rgb = color
        p.space_before = Pt(gap * 72 * 0.35)

def section_box(slide, title, l, t, w, h, col, items, item_size=12):
    rect(slide, l, t, w, 0.5, col)
    tx(slide, title, l+0.1, t+0.07, w-0.15, 0.4, size=13, bold=True, color=WHITE)
    rect(slide, l, t+0.5, w, h-0.5, WHITE)
    bullet_lines(slide, [(0, it) for it in items],
                 l+0.18, t+0.6, w-0.3, size=item_size, color=NAVY)

# ════════════════════════════════════════════════════════════
# S01  표지
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, NAVY)
rect(s, 0, 0, 13.33, 0.07, ACCENT)
rect(s, 0, 7.43, 13.33, 0.07, ACCENT)

tx(s, "반도체 파운드리 운영 대시보드", 0.8, 1.6, 11.73, 1.1,
   size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tx(s, "Semiconductor Foundry Operations Dashboard", 0.8, 2.65, 11.73, 0.6,
   size=19, color=RGBColor(0x90,0xaa,0xcc), align=PP_ALIGN.CENTER, italic=True)
rect(s, 4.2, 3.45, 4.93, 0.05, ACCENT)
tx(s, "AI 기반 생산·품질·매출 통합 분석 플랫폼  |  기존 시스템 연동 실운영 가이드",
   0.8, 3.6, 11.73, 0.55, size=14,
   color=RGBColor(0xcc,0xdd,0xee), align=PP_ALIGN.CENTER)

tags = ["Vue 3", "Spring Boot 3.2", "PostgreSQL 16",
        "Anthropic Claude API", "Java 21 LTS", "pgvector"]
widths  = [1.3, 1.9, 1.9, 2.5, 1.6, 1.5]
total_w = sum(widths) + (len(widths)-1)*0.15
start_x = (13.33 - total_w) / 2
cx = start_x
for tag, w in zip(tags, widths):
    badge(s, tag, cx, 4.45, w, 0.42, bg=BLUE, size=12)
    cx += w + 0.15

tx(s, "2026년 7월  |  개발팀  |  Ver 2.1  (RAG 문서 검색 확장 포함)",
   0.8, 6.6, 11.73, 0.45, size=13,
   color=RGBColor(0x77,0x88,0x99), align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# S02  목차
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, LIGHT)
header(s, "목차", "Agenda")

toc = [
    ("01", "프로젝트 배경 및 목표",         ACCENT),
    ("02", "시스템 아키텍처",               BLUE),
    ("03", "기존 시스템 연동 전략",          TEAL),
    ("04", "기술 스택",                     NAVY),
    ("05", "PostgreSQL 스키마 설계",        GREEN),
    ("06", "PostgreSQL 설치 및 권한 설정",  RGBColor(0x1b,0x6c,0x3a)),
    ("07", "AI 파이프라인 & 보안 가드레일", PURPLE),
    ("08", "프론트엔드 & REST API",         BLUE),
    ("09", "보안 가이드라인",               ACCENT),
    ("10", "배포 가이드 (Nginx / systemd)", TEAL),
    ("11", "주요 KPI 및 성과",              NAVY),
    ("13", "RAG 문서 검색 확장 (신규)",     GOLD),
    ("14", "향후 로드맵",                   GREEN),
]

cols = [toc[:7], toc[7:]]
for ci, col_items in enumerate(cols):
    lx = 0.45 + ci * 6.5
    for ri, (num, title, col) in enumerate(col_items):
        ty = 1.2 + ri * 0.86
        rect(s, lx, ty, 0.65, 0.72, col)
        tx(s, num, lx, ty+0.1, 0.65, 0.52,
           size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        rect(s, lx+0.65, ty, 5.5, 0.72, WHITE)
        tx(s, title, lx+0.8, ty+0.17, 5.2, 0.45, size=14, color=NAVY)

# ════════════════════════════════════════════════════════════
# S03  배경 및 목표
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, LIGHT)
header(s, "01  프로젝트 배경 및 목표", "Background & Objectives")

section_box(s, "현재 문제점 (Pain Points)", 0.35, 1.28, 5.95, 5.65, ACCENT, [
    "▪  생산·품질·매출 데이터가 MES/ERP/LIMS에 분산",
    "▪  데이터 조회 시 SQL 전문가 개입 필요 → 의사결정 지연",
    "▪  실시간 수율·불량 현황 즉시 파악 불가",
    "▪  임원 리포트를 수동으로 취합 → 주 1회 업데이트 한계",
    "▪  비개발자는 데이터 직접 분석 불가",
], item_size=13)

section_box(s, "기대 효과 (Expected Outcomes)", 6.55, 1.28, 6.45, 5.65, GREEN, [
    "✓  단일 대시보드에서 KPI 실시간 조회",
    "✓  AI 자연어 질의 → 누구나 데이터 분석 가능",
    "✓  SQL 가드레일 → DB 보안 유지",
    "✓  수율/불량 시각화 → 품질 이슈 즉시 감지",
    "✓  기존 MES/ERP 연동 → 별도 데이터 입력 불필요",
    "✓  향후 로컬 LLM 전환 → 데이터 완전 내부 보관",
], item_size=13)

# ════════════════════════════════════════════════════════════
# S04  시스템 아키텍처
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, LIGHT)
header(s, "02  시스템 아키텍처 (실운영 구성)", "Production System Architecture")

# Nginx
rect(s, 0.3, 1.25, 12.73, 0.72, RGBColor(0x22,0x22,0x40))
tx(s, "Nginx 리버스 프록시  |  /  → Vue 정적 파일  |  /api/*  → Spring Boot :8080  |  SSL/TLS 종단",
   0.5, 1.38, 12.3, 0.46, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 3 Layer
layers = [
    (ACCENT, "Frontend", "Vue 3 SPA\n(Nginx 정적 서빙)",
     ["KPI Cards × 4", "Chart.js 차트 × 4", "AI 채팅(출처 표시)+문서관리", "📎 업로드 버튼 (RAG)"]),
    (BLUE,   "Backend",  "Spring Boot 3.2\n:8080",
     ["REST API 8개 + /api/ai/chat", "AiChatService (SQL↔RAG 분류)", "SQL 가드레일 6종", "rag-service :8081 (RAG 확장)"]),
    (NAVY,   "Data",     "PostgreSQL 16\n:5432",
     ["5개 테이블 + 문서 2종(RAG)", "pgvector 확장 (임베딩 검색)", "6개 집계 View", "인덱스 14개 + ivfflat"]),
]
for i, (col, title, sub, items) in enumerate(layers):
    lx = 0.3 + i * 4.2
    rect(s, lx, 2.15, 3.85, 0.62, col)
    tx(s, title, lx+0.1, 2.18, 3.65, 0.35, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tx(s, sub,   lx+0.1, 2.48, 3.65, 0.38, size=10, color=WHITE, align=PP_ALIGN.CENTER, italic=True)
    rect(s, lx, 2.77, 3.85, 2.1, WHITE)
    for j, item in enumerate(items):
        tx(s, f"▸  {item}", lx+0.2, 2.85+j*0.47, 3.5, 0.42, size=11, color=NAVY)
    if i < 2:
        tx(s, "▶", lx+3.87, 3.6, 0.35, 0.5, size=18, color=GRAY, align=PP_ALIGN.CENTER)

# Claude API + Local LLM 전환 안내
rect(s, 0.3, 5.05, 8.1, 0.75, GOLD)
tx(s, "Anthropic Claude API  (현재)\nclaude-sonnet-4-6  |  Text-to-SQL + 한국어 요약",
   0.45, 5.1, 7.8, 0.65, size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
rect(s, 8.6, 5.05, 4.43, 0.75, RGBColor(0x33,0x55,0x33))
tx(s, "로컬 LLM 전환 가능\nvLLM + Qwen2.5-14B\n(외부망 차단 완전 내부 운영)",
   8.7, 5.1, 4.2, 0.65, size=10, color=WHITE, align=PP_ALIGN.CENTER)
tx(s, "↔ 전환", 8.35, 5.25, 0.3, 0.4, size=9, color=GRAY, align=PP_ALIGN.CENTER)

# 연동 화살표
rect(s, 3.35, 5.6, 5.4, 0.04, ACCENT)
tx(s, "▲  app.yml 1줄 변경만으로 전환", 3.5, 5.67, 5.0, 0.35, size=10, color=ACCENT)

# MES/ERP 연동
rect(s, 0.3, 6.1, 12.73, 0.72, RGBColor(0x22,0x33,0x44))
tx(s, "기존 운영 시스템 연동  |  MES (생산 Lot·수율)  |  ERP·SAP SD (수주·매출)  |  LIMS (불량 검사)  |  Option A/B/C 선택 가능",
   0.5, 6.23, 12.3, 0.46, size=11, color=RGBColor(0xaa,0xcc,0xdd), align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# S05  기존 시스템 연동 전략
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, LIGHT)
header(s, "03  기존 시스템 연동 전략", "Integration Strategy with Existing Systems")

options = [
    ("Option A", "직접 연결\n(Read-Only)", BLUE,
     ["기존 MES/ERP DB에 Read-Only 계정 생성",
      "Spring Boot → 기존 DB 직접 쿼리",
      "View로 스키마 정규화 레이어 구성",
      "✅ 실시간 데이터  ✅ 구현 단순",
      "⚠️ 기존 DB 부하 증가 가능",
      "→ 권장: 읽기 트래픽 낮은 환경"]),
    ("Option B", "ETL 배치\n(권장)", GREEN,
     ["별도 foundry_db 구성",
      "배치 스케줄러가 주기적으로 적재",
      "증분 ETL (updated_at 기준)",
      "✅ 기존 DB 무부하  ✅ 독립 운영",
      "⚠️ 데이터 최신성: 배치 주기 의존",
      "→ 권장: 초기 도입 표준 방식"]),
    ("Option C", "CDC 실시간\n(고급)", PURPLE,
     ["Debezium → Kafka → foundry_db",
      "DB 변경 즉시 스트리밍 동기화",
      "✅ 완전 실시간  ✅ 원본 DB 무부하",
      "⚠️ 구축 비용 높음 (Kafka 필요)",
      "⚠️ 운영 복잡도 증가",
      "→ 권장: 대규모 / 실시간 필수 환경"]),
]
for i, (opt, sub, col, items) in enumerate(options):
    lx = 0.35 + i * 4.3
    rect(s, lx, 1.28, 3.95, 0.65, col)
    tx(s, opt, lx+0.1, 1.3, 2.0, 0.35, size=17, bold=True, color=WHITE)
    tx(s, sub, lx+2.1, 1.32, 1.8, 0.58, size=11, color=WHITE, align=PP_ALIGN.RIGHT)
    rect(s, lx, 1.93, 3.95, 4.55, WHITE)
    for j, item in enumerate(items):
        c = GREEN if item.startswith("✅") else (ACCENT if item.startswith("⚠️") else
            (BLUE if item.startswith("→") else NAVY))
        tx(s, item, lx+0.18, 2.02+j*0.7, 3.65, 0.62, size=12, color=c)

# 단계별 전환 안내
rect(s, 0.35, 6.6, 12.63, 0.62, NAVY)
tx(s, "전환 로드맵:  초기 도입 → Option B (ETL 배치)  →  안정화 후 Option A 또는 C 전환  |  코드 변경 없이 application.yml 수정만으로 DB 소스 교체 가능",
   0.5, 6.7, 12.3, 0.44, size=11, color=WHITE, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# S06  기술 스택
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, LIGHT)
header(s, "04  기술 스택", "Technology Stack")

cols_data = [
    ("Frontend", ACCENT, [
        ("Vue 3.4",           "Composition API, 반응형 상태관리"),
        ("Vite 5.2",          "빠른 HMR, SPA 빌드"),
        ("Chart.js 4.4",      "Line / Bar / Doughnut 차트"),
        ("vue-chartjs 5.3",   "Vue 3 래퍼"),
        ("Axios 1.6",         "HTTP 클라이언트, 프록시"),
    ]),
    ("Backend", BLUE, [
        ("Spring Boot 3.2.5", "엔터프라이즈 REST API"),
        ("JdbcTemplate",      "경량 DB 접근, SQL 직접 제어"),
        ("HikariCP",          "고성능 커넥션 풀 (기본 포함)"),
        ("OpenJDK 21 LTS",    "최신 LTS, 가상 스레드 지원"),
        ("Maven 3.9",         "표준 Java 빌드 도구"),
    ]),
    ("Data & AI", NAVY, [
        ("PostgreSQL 16",     "RDBMS, 파티셔닝, FDW 지원"),
        ("Anthropic SDK 2.34","claude-sonnet-4-6 연동"),
        ("Regex 가드레일",    "DDL/DML/위험명령 6종 차단"),
        ("Nginx",             "리버스 프록시, 정적 서빙"),
        ("systemd",           "운영 서비스 관리"),
    ]),
]
for i, (cat, col, items) in enumerate(cols_data):
    lx = 0.35 + i * 4.3
    rect(s, lx, 1.28, 3.95, 0.55, col)
    tx(s, cat, lx+0.12, 1.32, 3.7, 0.48, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for j, (name, desc) in enumerate(items):
        bg = WHITE if j % 2 == 0 else LIGHT
        rect(s, lx, 1.83+j*1.06, 3.95, 1.0, bg)
        tx(s, name, lx+0.15, 1.88+j*1.06, 3.7, 0.42, size=13, bold=True, color=col)
        tx(s, desc, lx+0.15, 2.28+j*1.06, 3.7, 0.45, size=11, color=GRAY)

# ════════════════════════════════════════════════════════════
# S07  PostgreSQL 스키마 설계
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, LIGHT)
header(s, "05  PostgreSQL 스키마 설계", "Database Schema Design")

tables = [
    ("customers",        ACCENT, ["id, code (UNIQUE)", "name, short_name", "country, country_code", "tier  PLATINUM~BRONZE", "credit_limit, is_active"]),
    ("products",         BLUE,   ["id, code (UNIQUE)", "name, category", "technology_node", "wafer_size  150/200/300mm", "unit_price, cycle_time_days"]),
    ("sales_orders",     NAVY,   ["id, order_number (UNIQUE)", "customer_id FK / product_id FK", "quantity, unit_price", "total_amount (자동계산)", "status 9종 / priority 1~5"]),
    ("production_lots",  GREEN,  ["id, lot_number (UNIQUE)", "order_id FK / product_id FK", "quantity, start_date", "status 6종 / yield_rate", "fab_line, current_step"]),
    ("defect_records",   PURPLE, ["id, lot_id FK", "defect_type, defect_code", "process_step, equipment_id", "count, density, severity 4종", "root_cause, is_recurring"]),
]

# 위치 배치
pos = [(0.3,1.28),(2.75,1.28),(5.2,1.28),(0.9,4.1),(5.9,4.1)]
for (lx,ty),(name,col,fields) in zip(pos, tables):
    w = 2.3
    rect(s, lx, ty, w, 0.48, col)
    tx(s, name, lx+0.08, ty+0.06, w-0.12, 0.38, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, lx, ty+0.48, w, len(fields)*0.4+0.12, WHITE)
    for k, f in enumerate(fields):
        tx(s, f, lx+0.12, ty+0.54+k*0.4, w-0.2, 0.38, size=10, color=NAVY)

# FK 관계 텍스트
fk_lines = [
    "customers  1:N  sales_orders",
    "products   1:N  sales_orders",
    "products   1:N  production_lots",
    "production_lots  1:N  defect_records",
    "sales_orders  1:1  production_lots (order_id)",
]
rect(s, 8.7, 1.28, 4.3, 5.55, WHITE)
rect(s, 8.7, 1.28, 4.3, 0.48, TEAL)
tx(s, "테이블 관계 (FK)", 8.82, 1.3, 4.05, 0.42, size=13, bold=True, color=WHITE)
for k, line in enumerate(fk_lines):
    tx(s, f"▸  {line}", 8.85, 1.86+k*0.58, 4.05, 0.5, size=11, color=NAVY)

rect(s, 8.7, 4.7, 4.3, 2.13, LIGHT)
tx(s, "View (집계 레이어)", 8.85, 4.75, 4.05, 0.38, size=12, bold=True, color=NAVY)
views = ["v_kpi_summary", "v_revenue_monthly",
         "v_revenue_by_customer", "v_yield_by_product",
         "v_active_lots", "v_defect_summary"]
for k, v in enumerate(views):
    tx(s, f"◆  {v}", 8.9, 5.18+k*0.42, 4.0, 0.38, size=10, color=TEAL)

rect(s, 0.3, 6.88, 12.73, 0.5, NAVY)
tx(s, "total_amount 자동계산 컬럼 (GENERATED ALWAYS AS STORED)  |  updated_at 자동갱신 트리거  |  파티셔닝: 연간 1,000만 건 초과 시 권장",
   0.45, 6.95, 12.5, 0.38, size=10, color=WHITE, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# S08  PostgreSQL 설치 및 권한 설정
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, LIGHT)
header(s, "06  PostgreSQL 설치 및 권한 설정", "PostgreSQL Installation & Access Control")

# 왼쪽: 설치 방법
section_box(s, "OS별 설치 명령", 0.3, 1.28, 4.1, 3.5, BLUE, [
    "macOS:  brew install postgresql@16",
    "Ubuntu: apt install postgresql-16",
    "RHEL:   dnf install postgresql16-server",
    "Docker: postgres:16-alpine 이미지",
    "",
    "초기화  →  postgresql.conf 수정",
    "→  pg_hba.conf 설정  →  서비스 시작",
], item_size=11)

# 오른쪽: pg_hba.conf
rect(s, 4.6, 1.28, 8.43, 3.5, WHITE)
rect(s, 4.6, 1.28, 8.43, 0.48, NAVY)
tx(s, "pg_hba.conf — 접속 제어 (IP 화이트리스트)", 4.72, 1.3, 8.2, 0.42, size=13, bold=True, color=WHITE)
hba = [
    ("local", "all",        "postgres",    "-",                   "peer"),
    ("host",  "foundry_db", "foundry_app", "192.168.1.100/32",    "scram-sha-256"),
    ("host",  "foundry_db", "foundry_etl", "192.168.1.200/32",    "scram-sha-256"),
    ("host",  "foundry_db", "foundry_ro",  "192.168.10.0/24",     "scram-sha-256"),
    ("host",  "all",        "all",         "0.0.0.0/0",           "reject"),
]
hdrs = ["TYPE","DATABASE","USER","ADDRESS","METHOD"]
hw   = [0.75, 1.45, 1.35, 2.0, 1.7]
for ci,(h,w) in enumerate(zip(hdrs,hw)):
    lx = 4.72 + sum(hw[:ci])
    rect(s, lx, 1.82, w-0.03, 0.38, LIGHT)
    tx(s, h, lx+0.04, 1.84, w-0.08, 0.32, size=10, bold=True, color=GRAY)
for ri, row in enumerate(hba):
    bg = WHITE if ri%2==0 else RGBColor(0xf8,0xf8,0xff)
    mc = ACCENT if row[4]=="reject" else (GREEN if row[4]=="scram-sha-256" else BLUE)
    for ci,(cell,w) in enumerate(zip(row,hw)):
        lx = 4.72 + sum(hw[:ci])
        rect(s, lx, 2.25+ri*0.48, w-0.03, 0.44, bg)
        fc = mc if ci==4 else NAVY
        tx(s, cell, lx+0.04, 2.3+ri*0.48, w-0.08, 0.34, size=10, color=fc,
           bold=(ci==4))

# 아래: 계정 권한 분리
rect(s, 0.3, 4.93, 12.73, 0.48, NAVY)
tx(s, "계정 권한 분리 원칙 (최소 권한 원칙)", 0.45, 4.97, 12.5, 0.38, size=13, bold=True, color=WHITE)

accts = [
    ("foundry_app",  BLUE,   "SELECT만\n(대시보드 API)",    "max 20 connection"),
    ("foundry_ro",   TEAL,   "SELECT만\n(BI / 분석용)",     "max 10 connection"),
    ("foundry_etl",  GREEN,  "SELECT\nINSERT UPDATE DELETE", "max 5 connection\n(ETL 배치)"),
    ("postgres",     ACCENT, "슈퍼유저\n(DBA 전용)",        "앱 미사용\n로컬 소켓만"),
]
for i, (name, col, priv, note) in enumerate(accts):
    lx = 0.3 + i * 3.2
    rect(s, lx, 5.55, 3.0, 1.7, col)
    tx(s, name, lx+0.1, 5.6, 2.8, 0.4, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tx(s, priv, lx+0.1, 6.05, 2.8, 0.65, size=11, color=WHITE, align=PP_ALIGN.CENTER)
    tx(s, note, lx+0.1, 6.72, 2.8, 0.45, size=10, color=RGBColor(0xdd,0xee,0xff), align=PP_ALIGN.CENTER, italic=True)

# ════════════════════════════════════════════════════════════
# S09  AI 파이프라인
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, LIGHT)
header(s, "07  AI 파이프라인 & SQL 보안 가드레일", "AI Pipeline & Security Guardrail")

steps = [
    (ACCENT, "①\n사용자 질문\n(자연어)"),
    (BLUE,   "②\nClaude API\nSQL 생성"),
    (PURPLE, "③\n가드레일\n6종 검증"),
    (TEAL,   "④\nLIMIT 100\n자동 적용"),
    (NAVY,   "⑤\nPostgreSQL\n실행"),
    (GREEN,  "⑥\nClaude API\n한국어 요약"),
]
for i,(col,label) in enumerate(steps):
    lx = 0.3 + i*2.12
    rect(s, lx, 1.28, 1.85, 1.6, col)
    tx(s, label, lx+0.07, 1.33, 1.72, 1.48,
       size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(steps)-1:
        tx(s, "▶", lx+1.87, 1.85, 0.27, 0.5, size=16, color=GRAY, align=PP_ALIGN.CENTER)

# 가드레일 6종
rect(s, 0.3, 3.05, 12.73, 0.48, PURPLE)
tx(s, "가드레일 6종 규칙 상세 (Java Regex 기반)", 0.45, 3.08, 12.5, 0.42, size=14, bold=True, color=WHITE)

rules = [
    ("GR-01", "DDL 차단",       "CREATE · DROP · ALTER · TRUNCATE · RENAME",      ACCENT),
    ("GR-02", "DML 차단",       "INSERT · UPDATE · DELETE · MERGE · UPSERT",      ACCENT),
    ("GR-03", "위험명령 차단",   "GRANT · REVOKE · EXECUTE · EXEC · COPY · LOAD", ACCENT),
    ("GR-04", "SELECT 강제",    "쿼리가 SELECT로 시작하지 않으면 즉시 거부",         BLUE),
    ("GR-05", "LIMIT 자동화",   "없으면 100 자동 추가, 100 초과 시 100으로 교체",   GREEN),
    ("GR-06", "복수문장 차단",  "세미콜론(;) 이후 내용 존재 시 즉시 거부",           BLUE),
]
for i, (code, name, desc, col) in enumerate(rules):
    row = i // 2; ci = i % 2
    lx = 0.35 + ci * 6.5; ty = 3.67 + row * 0.95
    rect(s, lx, ty, 0.75, 0.75, col)
    tx(s, code, lx, ty+0.17, 0.75, 0.42, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, lx+0.75, ty, 5.55, 0.75, WHITE)
    tx(s, name, lx+0.85, ty+0.07, 1.6, 0.35, size=12, bold=True, color=col)
    tx(s, desc, lx+0.85, ty+0.38, 5.25, 0.32, size=11, color=NAVY)

rect(s, 0.3, 6.5, 12.73, 0.75, GOLD)
tx(s, "마크다운 코드블록 자동 제거: AI 응답에서 ```sql ... ``` 패턴을 제거 후 검증 적용\n모델 선택:  Haiku(빠름/저비용)  →  Sonnet(권장)  →  Opus(고정확도/고비용)",
   0.45, 6.55, 12.5, 0.65, size=11, color=NAVY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# S10  프론트엔드 & API
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, LIGHT)
header(s, "08  프론트엔드 & REST API 구성", "Frontend & REST API")

# 왼쪽: 컴포넌트
section_box(s, "Vue 3 컴포넌트 구조", 0.3, 1.28, 5.2, 5.65, ACCENT, [
    "App.vue",
    "  ▸  KPI Cards × 4 (매출/주문/수율/Lot)",
    "  ▸  Line Chart — 월별 매출 추이",
    "  ▸  Bar Chart (수평) — 고객별 매출",
    "  ▸  Bar Chart — 제품별 수율",
    "  ▸  Doughnut — 불량 유형 분포",
    "  ▸  생산 Lot 테이블 (최근 20건)",
    "  ▸  수주 테이블 (최근 20건)",
    "AIChatbox.vue — AI 자연어 채팅",
    "useDashboardAPI.js — Composable",
    "  ▸  Promise.all 병렬 API 호출 (7개)",
], item_size=12)

# 오른쪽: API 표
rect(s, 5.75, 1.28, 7.25, 5.65, WHITE)
rect(s, 5.75, 1.28, 7.25, 0.48, BLUE)
tx(s, "REST API  (Base: /api)", 5.88, 1.3, 7.0, 0.42, size=13, bold=True, color=WHITE)

apis = [
    ("GET",  "/kpis",                "KPI 4개 (매출·주문·수율·Lot)",  "< 100ms"),
    ("GET",  "/revenue-by-customer", "고객별 매출 집계",               "< 200ms"),
    ("GET",  "/revenue-trend",       "월별 매출 추이",                 "< 200ms"),
    ("GET",  "/yield-by-product",    "제품별 평균 수율",               "< 200ms"),
    ("GET",  "/defects-by-type",     "불량 유형·심각도 집계",          "< 200ms"),
    ("GET",  "/production-lots",     "최근 생산 Lot 20건",             "< 150ms"),
    ("GET",  "/orders",              "최근 수주 20건",                 "< 150ms"),
    ("POST", "/chat",                "AI 자연어 질의 → 요약",          "3~10s"),
]
hdrs2 = ["Method", "Endpoint", "설명", "목표"]
hw2   = [0.9, 2.4, 2.35, 1.3]
for ci,(h,w) in enumerate(zip(hdrs2,hw2)):
    lx = 5.88 + sum(hw2[:ci])
    rect(s, lx, 1.82, w-0.04, 0.38, LIGHT)
    tx(s, h, lx+0.05, 1.84, w-0.1, 0.32, size=10, bold=True, color=GRAY)
for ri, (method, path, desc, tgt) in enumerate(apis):
    bg = WHITE if ri%2==0 else RGBColor(0xf8,0xf9,0xff)
    mc = BLUE if method=="GET" else ACCENT
    row = (method, path, desc, tgt)
    for ci,(cell,w) in enumerate(zip(row,hw2)):
        lx = 5.88 + sum(hw2[:ci])
        rect(s, lx, 2.26+ri*0.57, w-0.04, 0.52, bg)
        fc = WHITE if ci==0 else (NAVY if ci<3 else GREEN)
        if ci==0:
            rect(s, lx, 2.26+ri*0.57, w-0.04, 0.52, mc)
        tx(s, cell, lx+0.05, 2.31+ri*0.57, w-0.1, 0.38,
           size=10, color=fc, bold=(ci==0))

# ════════════════════════════════════════════════════════════
# S11  보안 가이드라인
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, LIGHT)
header(s, "09  보안 가이드라인", "Security Guidelines")

sec_items = [
    ("DB 계정\n권한 분리", BLUE, [
        "foundry_app:  SELECT만",
        "foundry_ro:   SELECT만",
        "foundry_etl:  DML 허용",
        "postgres:     로컬 소켓만",
        "pg_hba.conf IP 화이트리스트",
    ]),
    ("네트워크\n방화벽", NAVY, [
        "PostgreSQL :5432",
        "→ 백엔드 IP만 허용",
        "API :8080",
        "→ Nginx IP만 허용",
        "외부 직접 접근 차단",
    ]),
    ("API Key\n관리", ACCENT, [
        ".env 파일 (.gitignore)",
        "운영: 환경변수 주입",
        "systemd EnvironmentFile",
        "Vault/Secrets Manager 권장",
        "코드에 하드코딩 절대 금지",
    ]),
    ("SQL\n인젝션 방지", GREEN, [
        "파라미터화 쿼리 사용",
        "JdbcTemplate ?  바인딩",
        "AI 생성 SQL 가드레일 검증",
        "사용자 입력 직접 조합 금지",
        "LIMIT 100 강제 적용",
    ]),
]
for i, (title, col, items) in enumerate(sec_items):
    lx = 0.3 + i * 3.2
    rect(s, lx, 1.28, 3.0, 0.75, col)
    tx(s, title, lx+0.1, 1.3, 2.8, 0.7, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, lx, 2.03, 3.0, 3.3, WHITE)
    for j, item in enumerate(items):
        tx(s, f"▸  {item}", lx+0.15, 2.12+j*0.6, 2.75, 0.52, size=11, color=NAVY)

# 환경별 설정
rect(s, 0.3, 5.5, 12.73, 0.48, TEAL)
tx(s, "환경별 설정 분리", 0.45, 5.53, 12.5, 0.38, size=13, bold=True, color=WHITE)

envs = [
    ("개발 (dev)",      "localhost DB\nSQL 로그 전체 출력\npool-size: 5"),
    ("스테이징 (stg)",  "내부망 DB\nINFO 로그\npool-size: 10"),
    ("운영 (prod)",     "운영 DB (환경변수)\nWARN 로그 + 파일 저장\npool-size: 20"),
]
for i, (env, desc) in enumerate(envs):
    lx = 0.3 + i * 4.25
    rect(s, lx, 6.1, 4.0, 1.2, WHITE)
    tx(s, env, lx+0.1, 6.15, 3.8, 0.4, size=13, bold=True, color=TEAL)
    tx(s, desc, lx+0.1, 6.55, 3.8, 0.65, size=11, color=NAVY)

# ════════════════════════════════════════════════════════════
# S12  배포 가이드
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, LIGHT)
header(s, "10  배포 가이드", "Deployment Guide")

# 왼쪽: Nginx
section_box(s, "Nginx 설정 요점", 0.3, 1.28, 5.95, 5.65, TEAL, [
    "/  →  Vue 3 dist/ 정적 파일 서빙",
    "try_files $uri /index.html  (SPA fallback)",
    "/api/  →  proxy_pass :8080",
    "/api/chat  →  proxy_read_timeout 120s",
    "SSL/TLS:  TLSv1.2 + TLSv1.3",
    "HTTP → HTTPS 리디렉트",
    "Cache-Control: public, 1h (정적파일)",
], item_size=12)

# 오른쪽: systemd
section_box(s, "systemd 서비스 설정", 6.55, 1.28, 6.45, 5.65, NAVY, [
    "User=foundry  (전용 OS 계정)",
    "WorkingDirectory=/opt/foundry-dashboard",
    "JAVA_OPTS=-Xms512m -Xmx2g -XX:+UseG1GC",
    "SPRING_PROFILES_ACTIVE=prod",
    "DB_HOST / DB_PASSWORD  (환경변수)",
    "ANTHROPIC_API_KEY  (환경변수)",
    "Restart=on-failure  RestartSec=10",
    "",
    "systemctl enable foundry-dashboard",
    "systemctl start  foundry-dashboard",
    "journalctl -u foundry-dashboard -f",
], item_size=11)

# 배포 절차 타임라인
rect(s, 0.3, 7.0, 12.73, 0.32, ACCENT)
tx(s, "배포 순서:  ① DB 마이그레이션  →  ② JAR 빌드 (mvn package)  →  ③ systemd 재시작  →  ④ npm run build  →  ⑤ Nginx reload  →  ⑥ Health check",
   0.45, 7.04, 12.5, 0.25, size=9, color=WHITE, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# S13  KPI 및 성과
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, LIGHT)
header(s, "11  주요 KPI 및 성과 (데모 데이터 기준)", "Key Performance Indicators")

kpis = [
    ("$22.4M", "총 매출",     BLUE),
    ("11건",   "활성 주문",   ACCENT),
    ("96.39%", "평균 수율",   GREEN),
    ("6개",    "진행 중 Lot", NAVY),
]
for i,(val,lbl,col) in enumerate(kpis):
    lx = 0.3 + i*3.2
    rect(s, lx, 1.28, 3.0, 1.5, col)
    tx(s, val, lx+0.1, 1.38, 2.8, 0.85,
       size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tx(s, lbl, lx+0.1, 2.18, 2.8, 0.45,
       size=12, color=RGBColor(0xcc,0xdd,0xee), align=PP_ALIGN.CENTER)

# 고객 Top 5
section_box(s, "고객별 매출 Top 5", 0.3, 2.95, 6.1, 4.0, BLUE, [
    "1위  TSMC (PLATINUM)            $7,360,000",
    "2위  Samsung (PLATINUM)         $6,280,000",
    "3위  Qualcomm (GOLD)            $3,304,000",
    "4위  MediaTek (SILVER)          $2,552,000",
    "5위  Intel Foundry (GOLD)       $1,155,000",
], item_size=12)

# 수율 Top 5
section_box(s, "제품별 수율 Top 5", 6.6, 2.95, 6.4, 4.0, GREEN, [
    "1위  Power IC 180nm             99.1 %",
    "2위  Logic 12nm                 98.0 %",
    "3위  DRAM 1z                    — (진행중)",
    "4위  Logic 7nm                  96.6 %",
    "5위  Logic 5nm                  95.5 %",
], item_size=12)

# ════════════════════════════════════════════════════════════
# S14  RAG 문서 검색 확장 (신규)
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, LIGHT)
header(s, "13  RAG 문서 검색 확장 (신규)", "Document RAG Extension — 2026-07-19 추가")

# 요약 바
rect(s, 0.3, 1.25, 12.73, 0.62, RGBColor(0x22,0x22,0x40))
tx(s, "사내 문서(규정·매뉴얼)를 업로드하면, 챗봇이 SQL 질문과 문서 질문을 자동으로 구분해 같은 창구에서 답변",
   0.5, 1.35, 12.3, 0.42, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 3 Layer (기존 아키텍처 슬라이드와 동일 패턴 재사용)
rag_layers = [
    (ACCENT, "Client", "Vue 3\n:5173",
     ["📎 문서 업로드 버튼", "답변 출처·신뢰도 표시", "문서 관리 화면 (목록/삭제)", "vue-router 없이 화면 토글"]),
    (BLUE,   "Backend",  "Spring Boot\n:8080 (기존)",
     ["질문 자동 분류 (SQL/문서)", "문서 API 프록시(중계)", "근거 기반 답변 생성(Claude)", "기존 NL2SQL 변경 없음"]),
    (GOLD,   "rag-service", "Spring Boot\n:8081 (신규)",
     ["Tika 파싱 (PDF/DOCX/TXT)", "청킹: 800자, 10% overlap", "임베딩 (OpenAI/해싱 폴백)", "pgvector 코사인 검색"]),
]
for i, (col, title, sub, items) in enumerate(rag_layers):
    lx = 0.3 + i * 4.2
    rect(s, lx, 2.1, 3.85, 0.62, col)
    fg = NAVY if col == GOLD else WHITE
    tx(s, title, lx+0.1, 2.13, 3.65, 0.35, size=15, bold=True, color=fg, align=PP_ALIGN.CENTER)
    tx(s, sub,   lx+0.1, 2.43, 3.65, 0.38, size=10, color=fg, align=PP_ALIGN.CENTER, italic=True)
    rect(s, lx, 2.72, 3.85, 2.05, WHITE)
    for j, item in enumerate(items):
        tx(s, f"▸  {item}", lx+0.2, 2.8+j*0.47, 3.5, 0.42, size=11, color=NAVY)
    if i < 2:
        tx(s, "▶", lx+3.87, 3.5, 0.35, 0.5, size=18, color=GRAY, align=PP_ALIGN.CENTER)

# 신규 API 4종
rect(s, 0.3, 4.95, 12.73, 0.45, PURPLE)
tx(s, "신규 API 4종  (모두 backend /api 경유, rag-service는 외부 미노출)", 0.45, 4.98, 12.5, 0.4,
   size=12, bold=True, color=WHITE)
apis4 = ["POST /api/ai/chat", "POST /api/documents/upload", "GET /api/documents", "DELETE /api/documents/{id}"]
aw = [2.85, 3.55, 2.7, 3.0]
cx = 0.45
for label, w in zip(apis4, aw):
    badge(s, label, cx, 5.5, w-0.15, 0.45, bg=NAVY, size=11)
    cx += w

# 왜 별도 모듈인가
rect(s, 0.3, 6.15, 12.73, 1.0, GOLD)
tx(s,
   "왜 별도 모듈(rag-service)인가:  검색(Retrieval)과 생성(Generation)의 책임을 분리 — "
   "rag-service는 검색만, backend는 기존처럼 Claude 호출/가드레일 전담\n"
   "왜 안전한가:  루트 pom.xml은 두 모듈을 묶기만 할 뿐 backend/pom.xml은 그대로 — "
   "기존 Docker/Railway 배포 파이프라인 무변경",
   0.5, 6.22, 12.3, 0.88, size=11.5, bold=True, color=NAVY)

# ════════════════════════════════════════════════════════════
# S15  향후 로드맵
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, LIGHT)
header(s, "14  향후 로드맵", "Future Roadmap")

roadmap = [
    ("0단계", "완료",  GOLD,   "RAG 문서 검색 (신규)",
     ["rag-service 모듈 + pgvector", "로컬 통합 검증 완료(07-19)", "→ 운영(Railway) 반영은 다음 단계"]),
    ("1단계", "2주",  ACCENT, "인증 & 권한 관리",
     ["JWT 기반 로그인", "역할별 화면 접근 제어", "감사 로그 기록"]),
    ("2단계", "4주",  BLUE,   "로컬 LLM 전환",
     ["vLLM + Qwen2.5-14B", "내부 GPU 서버 구축", "외부망 완전 차단 운영"]),
    ("3단계", "3주",  NAVY,   "실시간 MES 연동",
     ["CDC (Debezium+Kafka)", "WebSocket 실시간 차트", "수율 알림 경보"]),
]
for i,(stage,dur,col,title,items) in enumerate(roadmap):
    lx = 0.3 + i*3.2
    rect(s, lx, 1.28, 3.0, 0.68, col)
    tx(s, stage, lx+0.1, 1.3, 1.5, 0.35, size=16, bold=True, color=WHITE)
    tx(s, dur,   lx+1.6, 1.35, 1.3, 0.3, size=12, color=WHITE, align=PP_ALIGN.RIGHT)
    rect(s, lx, 1.96, 3.0, 0.48, RGBColor(0xee,0xee,0xff) if col==BLUE else LIGHT)
    tx(s, title, lx+0.12, 2.0, 2.8, 0.38, size=13, bold=True, color=col)
    rect(s, lx, 2.44, 3.0, 2.5, WHITE)
    for j, item in enumerate(items):
        tx(s, f"▸  {item}", lx+0.15, 2.52+j*0.72, 2.75, 0.62, size=12, color=NAVY)

# 로컬 LLM 핵심 강조
rect(s, 0.3, 5.1, 12.73, 1.35, GOLD)
tx(s, "핵심: 로컬 LLM 전환 (2단계)", 0.5, 5.15, 12.3, 0.45, size=15, bold=True, color=NAVY)
tx(s,
   "현재 Claude API  →  내부 GPU 서버 + Qwen2.5-14B (또는 Llama 3.3-70B)\n"
   "데이터가 외부로 전송되지 않아 사내 보안 정책 완전 준수  |  파운드리 공정 데이터 완전 내부 보관\n"
   "코드 변경 최소화:  application.yml URL·모델명 변경 + local-llm.enabled=true 로 전환 완료",
   0.5, 5.58, 12.3, 0.82, size=12, color=NAVY)

# ════════════════════════════════════════════════════════════
# S16  마무리
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, NAVY)
rect(s, 0, 0, 13.33, 0.07, ACCENT)
rect(s, 0, 7.43, 13.33, 0.07, ACCENT)

tx(s, "감사합니다", 0.8, 1.7, 11.73, 1.1,
   size=46, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tx(s, "Q & A", 0.8, 2.75, 11.73, 0.65,
   size=22, color=RGBColor(0x90,0xaa,0xcc), align=PP_ALIGN.CENTER)
rect(s, 4.0, 3.55, 5.33, 0.05, ACCENT)

# 요약 정보
info = [
    ("대시보드", "http://localhost:5173"),
    ("API 서버", "http://localhost:8080"),
    ("실행", "./start.sh"),
    ("문서", "docs/specification_enterprise.md"),
]
for i,(k,v) in enumerate(info):
    lx = 1.2 + (i%2)*5.8; ty = 4.1 + (i//2)*1.05
    tx(s, k, lx, ty, 2.5, 0.38, size=12, color=RGBColor(0x88,0x99,0xaa))
    tx(s, v, lx, ty+0.35, 5.5, 0.55, size=15, bold=True, color=WHITE)

# 문서 목록
rect(s, 2.0, 6.3, 9.33, 0.9, RGBColor(0x0f,0x1a,0x30))
tx(s,
   "산출물:  specification.md  |  specification_enterprise.md  |  postgresql_setup_guide.md  |  "
   "presentation.pptx  |  foundry-ai-proposal.pdf",
   2.1, 6.42, 9.1, 0.35, size=10, color=RGBColor(0xaa,0xbb,0xcc), align=PP_ALIGN.CENTER)

# ── 저장 ─────────────────────────────────────────────────────
out = "/Users/deejayseo/foundry-dashboard/docs/presentation.pptx"
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides)} slides)")
